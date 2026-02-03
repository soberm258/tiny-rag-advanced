from __future__ import annotations

from pathlib import Path
from typing import List


def read_pptx_to_text(path: Path) -> str:
    from pptx import Presentation

    ppt = Presentation(str(path))
    lines: List[str] = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)

